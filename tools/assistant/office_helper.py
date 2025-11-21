import io
import os
from typing import List, Optional

def create_text_file(path: str, text: str, dry_run: bool = True) -> str:
    if dry_run:
        return f"DRY_RUN: Would create text file at {path} with {len(text)} bytes"
    with open(path, "w", encoding="utf-8") as f:
        f.write(text)
    return f"Wrote text file {path}"

def create_docx(path: str, paragraphs: List[str], dry_run: bool = True) -> str:
    try:
        from docx import Document
    except Exception as e:
        return f"python-docx not installed: {e}"
    if dry_run:
        return f"DRY_RUN: Would create docx at {path} with {len(paragraphs)} paragraphs"
    doc = Document()
    for p in paragraphs:
        doc.add_paragraph(p)
    doc.save(path)
    return f"Wrote docx {path}"

def create_pptx(path: str, slides: List[dict], dry_run: bool = True) -> str:
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except Exception as e:
        return f"python-pptx not installed: {e}"
    if dry_run:
        return f"DRY_RUN: Would create pptx at {path} with {len(slides)} slides"
    prs = Presentation()
    for slide in slides:
        layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
        s = prs.slides.add_slide(layout)
        title = slide.get('title')
        body = slide.get('body')
        if title and s.shapes.title:
            s.shapes.title.text = title
        if body:
            try:
                tx = s.shapes.placeholders[1].text_frame
                tx.text = body
            except Exception:
                pass
    prs.save(path)
    return f"Wrote pptx {path}"
